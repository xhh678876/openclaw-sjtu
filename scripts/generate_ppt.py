#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
上海交通大学 PPT 生成工具
功能：基于模板生成 PPT，支持从 Markdown 自动生成多页幻灯片
依赖：python-pptx (pip install python-pptx)
模板目录：~/.openclaw/workspace/skills/sjtu-ppt/templates/
"""

import os
import sys
import re
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("❌ 缺少 python-pptx 库，请执行: pip install python-pptx")
    sys.exit(1)

SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATES_DIR = os.path.join(SKILL_DIR, "templates")
DEFAULT_TEMPLATE = os.path.join(TEMPLATES_DIR, "0.上海交通大学通用PPT模板.pptx")


def list_templates():
    """列出所有可用模板"""
    templates = []
    if not os.path.isdir(TEMPLATES_DIR):
        return templates
    for f in sorted(os.listdir(TEMPLATES_DIR)):
        if f.lower().endswith((".pptx", ".ppt")) and not f.startswith("~$"):
            fpath = os.path.join(TEMPLATES_DIR, f)
            size_kb = os.path.getsize(fpath) / 1024
            templates.append({
                "name": f,
                "path": fpath,
                "size": f"{size_kb:.1f} KB",
            })
    return templates


def _get_template_path(template_name=None):
    """解析模板路径"""
    if template_name is None:
        if os.path.exists(DEFAULT_TEMPLATE):
            return DEFAULT_TEMPLATE
        return None  # 无模板，创建空白 PPT

    # 精确路径
    if os.path.exists(template_name):
        return template_name

    # 在模板目录中搜索
    if os.path.isdir(TEMPLATES_DIR):
        for f in os.listdir(TEMPLATES_DIR):
            if template_name.lower() in f.lower() and f.lower().endswith(".pptx"):
                return os.path.join(TEMPLATES_DIR, f)

    return None


def _get_layout(prs, layout_name):
    """获取幻灯片版式，容错处理"""
    layout_map = {
        "title": 0,
        "content": 1,
        "section": 2,
        "blank": 6,
    }
    idx = layout_map.get(layout_name, 1)
    layouts = prs.slide_layouts
    if idx < len(layouts):
        return layouts[idx]
    # 回退：返回第一个可用版式
    return layouts[0] if layouts else None


def _add_text_to_placeholder(placeholder, text, font_size=18):
    """向占位符添加文本"""
    tf = placeholder.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line.strip()
        if p.runs:
            p.runs[0].font.size = Pt(font_size)


def generate_ppt(title, slides_content, template_path=None, output_path="output.pptx"):
    """
    生成 PPT

    参数:
        title: PPT 标题（用于首页）
        slides_content: 幻灯片内容列表，每项为 dict:
            {"title": "标题", "content": "正文", "layout": "title/content/section/blank"}
        template_path: 模板文件路径（None 则使用默认模板）
        output_path: 输出路径
    """
    tpl = _get_template_path(template_path)

    if tpl and os.path.exists(tpl):
        prs = Presentation(tpl)
    else:
        prs = Presentation()
        if tpl is not None:
            print(f"⚠️  模板未找到: {tpl}，使用空白模板")

    # 添加标题页
    title_layout = _get_layout(prs, "title")
    if title_layout:
        slide = prs.slides.add_slide(title_layout)
        if slide.placeholders:
            # 标题
            if 0 in slide.placeholders:
                slide.placeholders[0].text = title
            # 副标题
            if 1 in slide.placeholders:
                slide.placeholders[1].text = "上海交通大学"

    # 添加内容页
    for item in slides_content:
        layout_name = item.get("layout", "content")
        layout = _get_layout(prs, layout_name)
        if not layout:
            continue

        slide = prs.slides.add_slide(layout)
        placeholders = slide.placeholders

        slide_title = item.get("title", "")
        slide_content = item.get("content", "")

        if layout_name == "blank":
            # 空白页：如果有内容，添加文本框
            if slide_content:
                from pptx.util import Inches
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
                tf = txBox.text_frame
                tf.word_wrap = True
                for i, line in enumerate(slide_content.split("\n")):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line.strip()
        else:
            # 有占位符的版式
            if 0 in placeholders and slide_title:
                placeholders[0].text = slide_title
            if 1 in placeholders and slide_content:
                _add_text_to_placeholder(placeholders[1], slide_content)

    # 确保输出目录存在
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    return {"success": True, "path": os.path.abspath(output_path), "slides": len(slides_content) + 1}


def generate_from_markdown(title, markdown_text, template_path=None, output_path="output.pptx"):
    """
    从 Markdown 文本生成 PPT
    按 ## 标题切分为各页，支持 # 作为节标题

    参数:
        title: PPT 标题
        markdown_text: Markdown 文本内容
        template_path: 模板路径
        output_path: 输出路径
    """
    slides_content = []

    # 如果是文件路径，读取文件
    if os.path.isfile(markdown_text):
        with open(markdown_text, "r", encoding="utf-8") as f:
            markdown_text = f.read()

    # 按标题切分
    # 先按 ## 切分（二级标题作为幻灯片标题）
    sections = re.split(r'^(#{1,2})\s+(.+)$', markdown_text, flags=re.MULTILINE)

    if len(sections) <= 1:
        # 没有标题标记，整个内容作为一页
        slides_content.append({
            "title": title,
            "content": markdown_text.strip(),
            "layout": "content",
        })
    else:
        # 开头如果有内容（在第一个标题之前）
        preamble = sections[0].strip()
        if preamble:
            slides_content.append({
                "title": "概述",
                "content": _clean_markdown(preamble),
                "layout": "content",
            })

        # 解析标题和内容
        i = 1
        while i < len(sections) - 2:
            level = sections[i]      # # 或 ##
            heading = sections[i+1]  # 标题文本
            content = sections[i+2].strip() if i+2 < len(sections) else ""
            i += 3

            if level == "#":
                layout = "section"
            else:
                layout = "content"

            slides_content.append({
                "title": heading.strip(),
                "content": _clean_markdown(content),
                "layout": layout,
            })

    return generate_ppt(title, slides_content, template_path, output_path)


def _clean_markdown(text):
    """清理 Markdown 标记，转为纯文本"""
    # 移除图片
    text = re.sub(r'!\[.*?\]\(.*?\)', '', text)
    # 链接保留文本
    text = re.sub(r'\[(.+?)\]\(.*?\)', r'\1', text)
    # 移除加粗/斜体标记
    text = re.sub(r'\*{1,2}(.+?)\*{1,2}', r'\1', text)
    # 列表项保留
    text = re.sub(r'^[-*+]\s+', '• ', text, flags=re.MULTILINE)
    # 移除代码块标记
    text = re.sub(r'```\w*\n?', '', text)
    # 移除行内代码标记
    text = re.sub(r'`(.+?)`', r'\1', text)
    # 清理多余空行
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def main():
    parser = argparse.ArgumentParser(description="上海交通大学 PPT 生成工具")
    parser.add_argument("--title", "-t", help="PPT 标题")
    parser.add_argument("--markdown", "-m", help="Markdown 内容或文件路径")
    parser.add_argument("--template", help="模板名称或路径 (默认: 交大通用模板)")
    parser.add_argument("--output", "-o", default="output.pptx", help="输出文件路径")
    parser.add_argument("--list-templates", action="store_true", help="列出所有可用模板")
    args = parser.parse_args()

    if args.list_templates:
        templates = list_templates()
        if not templates:
            print(f"📂 模板目录为空: {TEMPLATES_DIR}")
        else:
            print(f"\n📂 可用模板 ({TEMPLATES_DIR})")
            print("=" * 50)
            for t in templates:
                print(f"  📄 {t['name']}  ({t['size']})")
        print()
        return

    if not args.title:
        print("❌ 错误: 请提供 --title 参数")
        sys.exit(1)

    if args.markdown:
        result = generate_from_markdown(
            title=args.title,
            markdown_text=args.markdown,
            template_path=args.template,
            output_path=args.output,
        )
    else:
        # 无 markdown 时创建只有标题页的 PPT
        result = generate_ppt(
            title=args.title,
            slides_content=[],
            template_path=args.template,
            output_path=args.output,
        )

    if result["success"]:
        print(f"✅ PPT 生成成功!")
        print(f"   📄 文件: {result['path']}")
        print(f"   📊 页数: {result['slides']}")
    else:
        print("❌ PPT 生成失败")
        sys.exit(1)


if __name__ == "__main__":
    main()
